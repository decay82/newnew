from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Define colors (Notion-inspired palette)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(55, 53, 47)  # Notion dark text
GRAY = RGBColor(120, 119, 116)  # Notion gray
ACCENT = RGBColor(46, 170, 220)  # Blue accent
LIGHT_BG = RGBColor(247, 246, 243)  # Notion light background

def add_blank_slide(prs):
    """Add a blank slide"""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(blank_layout)

def set_background_color(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    """Add a text box to slide"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True

    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.name = 'Arial'
            run.font.bold = bold
            run.font.color.rgb = color

    return textbox

# Slide 1: Cover/Home Page
slide1 = add_blank_slide(prs)
set_background_color(slide1, WHITE)

# Title
add_text_box(slide1, Inches(1), Inches(3), Inches(14), Inches(1),
             "🎮 게임 이름", font_size=72, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide1, Inches(1), Inches(4.5), Inches(14), Inches(0.8),
             "새로운 모험이 시작됩니다", font_size=28, color=GRAY, align=PP_ALIGN.CENTER)

# CTA Button (simulated with text box)
cta_box = slide1.shapes.add_shape(1, Inches(6.5), Inches(6), Inches(3), Inches(0.8))
cta_box.fill.solid()
cta_box.fill.fore_color.rgb = ACCENT
cta_box.line.color.rgb = ACCENT
add_text_box(slide1, Inches(6.5), Inches(6.1), Inches(3), Inches(0.6),
             "지금 플레이하기", font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# Slide 2: About the Game
slide2 = add_blank_slide(prs)
set_background_color(slide2, LIGHT_BG)

add_text_box(slide2, Inches(1), Inches(1), Inches(14), Inches(1),
             "게임 소개", font_size=48, bold=True, color=BLACK)

# Content boxes
content_y = 2.5
content_texts = [
    "🎯 장르: 액션 어드벤처 / RPG",
    "👥 플레이어: 싱글/멀티플레이 지원",
    "⏱️ 플레이 타임: 20-30시간",
    "🌍 오픈 월드 탐험과 다양한 퀘스트를 경험하세요"
]

for text in content_texts:
    add_text_box(slide2, Inches(2), Inches(content_y), Inches(12), Inches(0.7),
                 text, font_size=24, color=BLACK)
    content_y += 1


# Slide 3: Key Features
slide3 = add_blank_slide(prs)
set_background_color(slide3, WHITE)

add_text_box(slide3, Inches(1), Inches(0.8), Inches(14), Inches(1),
             "주요 특징", font_size=48, bold=True, color=BLACK)

# Feature cards
features = [
    ("🗺️", "광활한 오픈월드", "끝없이 펼쳐진 세계를 자유롭게 탐험하세요"),
    ("⚔️", "다이나믹한 전투", "실시간 전투 시스템으로 짜릿한 액션을 경험하세요"),
    ("🎨", "아름다운 그래픽", "차세대 그래픽 엔진으로 구현된 생생한 비주얼"),
    ("📖", "몰입감 있는 스토리", "플레이어의 선택이 결말을 바꾸는 스토리")
]

feature_x = 1
feature_y = 2.5
feature_width = 3.2
card_count = 0

for emoji, title, desc in features:
    # Feature box
    box = slide3.shapes.add_shape(1, Inches(feature_x), Inches(feature_y),
                                   Inches(feature_width), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.fill.background()

    # Emoji
    add_text_box(slide3, Inches(feature_x + 0.2), Inches(feature_y + 0.3),
                 Inches(feature_width - 0.4), Inches(0.6),
                 emoji, font_size=48, align=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide3, Inches(feature_x + 0.2), Inches(feature_y + 1.1),
                 Inches(feature_width - 0.4), Inches(0.5),
                 title, font_size=22, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide3, Inches(feature_x + 0.2), Inches(feature_y + 1.7),
                 Inches(feature_width - 0.4), Inches(0.7),
                 desc, font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

    card_count += 1
    if card_count == 2:
        feature_x = 1
        feature_y = 5.5
    else:
        feature_x += feature_width + 0.5


# Slide 4: Screenshots Gallery
slide4 = add_blank_slide(prs)
set_background_color(slide4, WHITE)

add_text_box(slide4, Inches(1), Inches(0.8), Inches(14), Inches(1),
             "스크린샷 갤러리", font_size=48, bold=True, color=BLACK)

# Screenshot placeholders (3 boxes)
screenshot_positions = [
    (1, 2.5, 4.5, 2.5),
    (5.7, 2.5, 4.5, 2.5),
    (10.4, 2.5, 4.5, 2.5),
    (1, 5.2, 4.5, 2.5),
    (5.7, 5.2, 4.5, 2.5),
    (10.4, 5.2, 4.5, 2.5)
]

for left, top, width, height in screenshot_positions:
    box = slide4.shapes.add_shape(1, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(230, 230, 230)
    box.line.color.rgb = RGBColor(200, 200, 200)

    # Placeholder text
    add_text_box(slide4, Inches(left), Inches(top + height/2 - 0.3),
                 Inches(width), Inches(0.6),
                 "게임 스크린샷", font_size=18, color=GRAY, align=PP_ALIGN.CENTER)


# Slide 5: Download/Get Started
slide5 = add_blank_slide(prs)
set_background_color(slide5, ACCENT)

add_text_box(slide5, Inches(1), Inches(2.5), Inches(14), Inches(1),
             "지금 바로 시작하세요!", font_size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide5, Inches(1), Inches(3.8), Inches(14), Inches(0.8),
             "다양한 플랫폼에서 즐길 수 있습니다", font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

# Platform buttons
platforms = ["🎮 Steam", "🎯 Epic Games", "📱 Mobile"]
btn_y = 5.5
btn_width = 3.5

for i, platform in enumerate(platforms):
    btn_x = 2.5 + (i * 4)

    # Button background
    btn = slide5.shapes.add_shape(1, Inches(btn_x), Inches(btn_y),
                                   Inches(btn_width), Inches(0.9))
    btn.fill.solid()
    btn.fill.fore_color.rgb = WHITE
    btn.line.fill.background()

    # Button text
    add_text_box(slide5, Inches(btn_x), Inches(btn_y + 0.15),
                 Inches(btn_width), Inches(0.6),
                 platform, font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# Slide 6: Contact/Community
slide6 = add_blank_slide(prs)
set_background_color(slide6, LIGHT_BG)

add_text_box(slide6, Inches(1), Inches(2), Inches(14), Inches(1),
             "커뮤니티에 참여하세요", font_size=48, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# Social media/contact info
contact_y = 4
contacts = [
    "🌐 웹사이트: www.yourgame.com",
    "💬 디스코드: discord.gg/yourgame",
    "📧 이메일: contact@yourgame.com",
    "🐦 트위터: @YourGame"
]

for contact in contacts:
    add_text_box(slide6, Inches(4), Inches(contact_y), Inches(8), Inches(0.6),
                 contact, font_size=22, color=BLACK, align=PP_ALIGN.CENTER)
    contact_y += 0.8

# Save presentation
output_file = '/home/user/newnew/게임_소개_사이트_디자인.pptx'
prs.save(output_file)
print(f"✅ PPT 파일이 생성되었습니다: {output_file}")
